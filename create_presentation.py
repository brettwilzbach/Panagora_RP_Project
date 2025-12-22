from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Panagora brand colors from website
NAVY = RGBColor(0, 51, 102)           # Deep blue from Panagora site
ROYAL_BLUE = RGBColor(0, 102, 153)    # Teal accent from Panagora
GOLD = RGBColor(0, 153, 153)          # Teal/cyan highlight
CHARCOAL = RGBColor(51, 51, 51)       # Body text
LIGHT_GRAY = RGBColor(240, 243, 245)  # Light background
WHITE = RGBColor(255, 255, 255)

LOGO_PATH = r"c:\Users\bwilzbach\Panagora\panagora_logo.png"

def add_title_slide(prs, title, subtitle, tagline=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.15))
    accent.fill.solid()
    accent.fill.fore_color.rgb = GOLD
    accent.line.fill.background()
    # Add logo
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(10.5), Inches(0.4), height=Inches(0.8))
    except:
        pass
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(2.2), Inches(11.733), Inches(1.8))
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    sub_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.0), Inches(11.733), Inches(0.8))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(22)
    p.font.color.rgb = GOLD
    p.alignment = PP_ALIGN.LEFT
    bottom = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(6.8), Inches(13.333), Inches(0.7))
    bottom.fill.solid()
    bottom.fill.fore_color.rgb = ROYAL_BLUE
    bottom.line.fill.background()
    if tagline:
        tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(6.9), Inches(11.733), Inches(0.5))
        tf = tag_box.text_frame
        p = tf.paragraphs[0]
        p.text = tagline
        p.font.size = Pt(16)
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.LEFT
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), Inches(13.333), Inches(0.06))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLD
    gold_line.line.fill.background()
    # Add logo
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(11.5), Inches(0.35), height=Inches(0.7))
    except:
        pass
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.5), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.9), Inches(5.9), Inches(5.0))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GRAY
    left_box.line.color.rgb = RGBColor(200, 200, 200)
    left_header = slide.shapes.add_textbox(Inches(0.7), Inches(2.1), Inches(5.5), Inches(0.5))
    tf = left_header.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = NAVY
    left_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.7), Inches(5.5), Inches(4.0))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = CHARCOAL
        p.space_before = Pt(10) if i > 0 else Pt(0)
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.9), Inches(5.9), Inches(5.0))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = NAVY
    right_box.line.fill.background()
    right_header = slide.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.5), Inches(0.5))
    tf = right_header.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = GOLD
    right_content = slide.shapes.add_textbox(Inches(7.1), Inches(2.7), Inches(5.5), Inches(4.0))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(14)
        p.font.color.rgb = WHITE
        p.space_before = Pt(10) if i > 0 else Pt(0)
    return slide

def add_content_slide_with_sections(prs, title, sections, personalized_for=None):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(1.5))
    header.fill.solid()
    header.fill.fore_color.rgb = NAVY
    header.line.fill.background()
    gold_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(1.5), Inches(13.333), Inches(0.06))
    gold_line.fill.solid()
    gold_line.fill.fore_color.rgb = GOLD
    gold_line.line.fill.background()
    # Add logo
    try:
        slide.shapes.add_picture(LOGO_PATH, Inches(11.5), Inches(0.35), height=Inches(0.7))
    except:
        pass
    # Add personalization text if provided
    if personalized_for:
        pers_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.1), Inches(10), Inches(0.3))
        tf = pers_box.text_frame
        p = tf.paragraphs[0]
        p.text = personalized_for
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = GOLD
        p.alignment = PP_ALIGN.LEFT
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(10.5), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.LEFT
    section_width = 3.9
    section_gap = 0.3
    start_x = 0.5
    for i, section in enumerate(sections):
        x_pos = start_x + i * (section_width + section_gap)
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(1.85), Inches(section_width), Inches(5.3))
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_GRAY
        box.line.color.rgb = RGBColor(200, 200, 200)
        circle = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x_pos + 0.15), Inches(2.0), Inches(0.45), Inches(0.45))
        circle.fill.solid()
        circle.fill.fore_color.rgb = ROYAL_BLUE
        circle.line.fill.background()
        num_box = slide.shapes.add_textbox(Inches(x_pos + 0.15), Inches(2.05), Inches(0.45), Inches(0.4))
        tf = num_box.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER
        header_box = slide.shapes.add_textbox(Inches(x_pos + 0.7), Inches(2.0), Inches(section_width - 0.9), Inches(0.6))
        tf = header_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = section['header']
        p.font.size = Pt(15)
        p.font.bold = True
        p.font.color.rgb = NAVY
        items_box = slide.shapes.add_textbox(Inches(x_pos + 0.2), Inches(2.65), Inches(section_width - 0.4), Inches(4.3))
        tf = items_box.text_frame
        tf.word_wrap = True
        for j, item in enumerate(section['items']):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = f"» {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = CHARCOAL
            p.space_before = Pt(8) if j > 0 else Pt(0)
    return slide

add_title_slide(prs, "Accelerating Growth\nin the AI Era", "Distribution Strategy | Investor Engagement | AUM Expansion", "Panagora Asset Management  |  Strategic Discussion  |  December 2024")

add_two_column_slide(prs, "The Quantitative Investing Landscape & Panagora's Edge", "MARKET DYNAMICS",
    ["» $2T+ in systematic AUM globally—quant now mainstream", "» Shift from 'big data' to 'smart data' with demonstrable alpha", "» ESG integration table stakes, but quality data scarce", "» Risk Parity gaining traction: 60/40 portfolios carry 93% equity risk", "» AI reshaping both investment process AND distribution expectations"],
    "PANAGORA DIFFERENTIATORS",
    ["» Proprietary data creation: biotech FDA models, NLP sentiment, Chinese social scraping", "» 'Discovery & Dollars' philosophy—human-machine synthesis", "» Integrated research + portfolio management teams", "» Risk Parity leadership: 0.87 vs 0.67 Sharpe (vs 60/40)", "» ESG innovation: greenwashing detection via NLP, materiality focus"])

add_content_slide_with_sections(prs, "AI-Era Distribution: Intelligent Lead Management & Engagement", [
    {'header': 'Predictive Lead Scoring', 'items': ["ML models on CRM data: score by conversion probability & AUM potential", "External signals: job changes, RFP activity, regulatory filings, conference attendance", "Real-time alerts on buying intent: website visits, content downloads", "Prioritize sales time on highest-value opportunities"]},
    {'header': 'Hyper-Personalized Outreach', 'items': ["AI-drafted communications tailored to portfolio gaps & interests", "Dynamic content: Risk Parity to consultants, ESG to pensions", "Adaptive nurture sequences based on engagement patterns", "Scale personalization without scaling headcount"]},
    {'header': 'Intelligent Meeting Prep', 'items': ["AI briefing docs: holdings, recent allocations, stated concerns", "Competitive positioning summaries before every pitch", "Conversation guides highlighting relevant Panagora strengths", "Post-meeting auto-generated follow-up recommendations"]}],
    personalized_for="Created for Discussion with Tim Stanton")

add_content_slide_with_sections(prs, "Generating Buzz: Thought Leadership & Publicity for Modern Investors", [
    {'header': 'AI-Powered Content Engine', 'items': ["Research → white papers → podcasts → social threads → video clips", "Real-time market commentary: AI draft + human polish = speed + quality", "Interactive 'Risk Parity Portfolio Builder' for prospects", "'Smart Data, Smart Alpha' thought leadership series"]},
    {'header': 'Strategic Amplification', 'items': ["LinkedIn optimization: timing, hashtags, engagement triggers", "Earned media: pitch 'Why 60/40 is dead' angles to financial press", "Podcast circuit: CIO shows with Eric's 'pilot + quant' narrative", "Update & amplify Institutional Investor Risk Parity research"]},
    {'header': 'Community & Events', 'items': ["Virtual 'Quant Masterclass' series—exclusive, invite-only", "AI-powered post-event follow-up within 24 hours", "Target ICP: diversification seekers, pension CIOs", "Build long-term relationships, not just transactions"]}])

prs.save(r"c:\Users\bwilzbach\Panagora\Panagora_Growth_Strategies.pptx")
print("Done!")
